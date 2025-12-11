import pandas as pd
from pptx import Presentation
from pptx.util import Pt
import os
import re

def generate_presentations_from_csv(csv_path: str, 
                                   template_path: str, 
                                   output_dir: str) -> int:
    """
    Generate PowerPoint presentations dari CSV hasil analisis
    
    Returns:
        int: Jumlah presentasi yang berhasil dibuat
    """
    
    print("\n" + "="*60)
    print("📊 GENERATING POWERPOINT PRESENTATIONS")
    print("="*60)
    
    # Read CSV/Excel - PERBAIKAN ENCODING
    try:
        print(f"Mencoba membaca file: {csv_path}")
        
        # Cek apakah file CSV atau Excel
        if csv_path.lower().endswith('.csv'):
            # Coba berbagai encoding
            encodings = ['utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
            
            df = None
            encoding_used = None
            
            for encoding in encodings:
                try:
                    print(f"  Mencoba encoding: {encoding}")
                    df = pd.read_csv(csv_path, encoding=encoding)
                    encoding_used = encoding
                    print(f"  ✓ Berhasil dengan encoding: {encoding}")
                    break
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    print(f"  ✗ Error dengan {encoding}: {e}")
                    continue
            
            if df is None:
                # Fallback: baca sebagai binary dan decode
                try:
                    with open(csv_path, 'rb') as f:
                        content = f.read()
                    
                    # Coba decode dengan replace errors
                    content_decoded = content.decode('utf-8', errors='replace')
                    
                    # Baca dari string
                    from io import StringIO
                    df = pd.read_csv(StringIO(content_decoded))
                    print(f"  ✓ Berhasil dengan binary read + replace errors")
                except Exception as e:
                    print(f"  ❌ Gagal semua encoding: {e}")
                    return 0
                    
        elif csv_path.lower().endswith(('.xlsx', '.xls')):
            # Baca dari Excel langsung
            try:
                df = pd.read_excel(csv_path)
                print(f"  ✓ Berhasil membaca Excel file")
            except Exception as e:
                print(f"  ❌ Error membaca Excel: {e}")
                return 0
        else:
            print(f"  ❌ Format file tidak didukung: {csv_path}")
            return 0
            
        print(f"✅ Loaded {len(df)} records")
        print(f"  Columns: {list(df.columns)}")
        
    except Exception as e:
        print(f"❌ Error reading file {csv_path}: {e}")
        print(f"   Error type: {type(e).__name__}")
        return 0
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    successful_count = 0
    
    for index, row in df.iterrows():
        try:
            # Dapatkan nama dengan berbagai cara
            nama = None
            
            # Coba berbagai kemungkinan nama kolom
            name_columns = ['nama', 'Nama', 'name', 'Name', 'candidate_name', 'full_name']
            
            for col in name_columns:
                if col in row and pd.notna(row[col]):
                    nama = str(row[col]).strip()
                    break
            
            if nama is None:
                nama = f"Candidate_{index + 1}"
                print(f"  ⚠ Nama tidak ditemukan, menggunakan: {nama}")
            
            print(f"\n[{index + 1}/{len(df)}] 📄 Generating for: {nama}")
            
            # Load template
            try:
                prs = Presentation(template_path)
                print(f"  ✓ Template loaded: {template_path}")
            except Exception as e:
                print(f"  ❌ Error loading template: {e}")
                continue
            
            # Cek slide
            if len(prs.slides) == 0:
                print(f"  ⚠ Template tidak memiliki slide")
                continue
            
            slide = prs.slides[0]
            
            # Debug: tampilkan data row
            print(f"  Data columns: {[col for col in row.index if pd.notna(row[col])]}")
            
            # Replace placeholders
            try:
                replace_placeholders(slide, row)
                print(f"  ✓ Placeholders replaced")
            except Exception as e:
                print(f"  ⚠ Error replacing placeholders: {e}")
                # Lanjutkan meskipun ada error
            
            # Save presentation
            try:
                # Bersihkan nama file dari karakter tidak valid
                clean_name = re.sub(r'[<>:"/\\|?*]', '_', nama)
                output_path = os.path.join(output_dir, f"Resume_{clean_name}.pptx")
                
                prs.save(output_path)
                print(f"  ✅ Saved: {os.path.basename(output_path)}")
                successful_count += 1
                
            except Exception as e:
                print(f"  ❌ Error saving presentation: {e}")
                continue
            
        except Exception as e:
            print(f"  ❌ Error generating for row {index}: {e}")
            import traceback
            traceback.print_exc()
            continue
    
    print(f"\n{'='*60}")
    if successful_count > 0:
        print(f"✅ Successfully generated {successful_count}/{len(df)} presentations")
        print(f"📁 Output folder: {output_dir}")
    else:
        print(f"❌ Failed to generate any presentations")
    
    return successful_count

def replace_placeholders(slide, row):
    """
    Replace placeholders dalam slide dengan data dari row
    """
    
    print(f"    Memproses placeholders untuk row...")
    print(f"    Row columns: {[col for col in row.index if pd.notna(row[col])]}")
    print(f"    Row data - nik: {row.get('nik', 'NOT FOUND')}")
    print(f"    Row data - nama: {row.get('nama', 'NOT FOUND')}")
    
    # Mapping placeholder dengan semua kemungkinan nama kolom
    placeholder_mappings = {
        '{{nik}}': ['nik', 'id', 'employee_id', 'employee id', 'no_induk', 'nomor induk'],
        '{{nama}}': ['nama', 'Nama', 'name', 'Name', 'candidate_name'],
        '{{executive summary}}': ['summary executive', 'summary_executive', 'executive_summary', 'summary'],
        '{{education}}': ['education', 'Education', 'pendidikan', 'Pendidikan'],
        '{{jabatan terakhir}}': ['jabatan terakhir', 'jabatan', 'position', 'jabatan_terakhir', 'current_position'],
        '{{competency}}': ['competency', 'Competency', 'skills', 'Skills', 'competency_data'],
        '{{experience}}': ['experience', 'Experience', 'pengalaman', 'pengalaman_kerja'],
        '{{business impact}}': ['business impact', 'business_impact', 'impact', 'business_impact_data']
    }
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        original_text = text_frame.text
        
        if not original_text:
            continue
        
        print(f"      Original text: {original_text[:100]}...")
        
        new_text = original_text
        
        # Cek setiap placeholder pattern
        for placeholder, possible_columns in placeholder_mappings.items():
            # Cari placeholder dengan flexible matching
            if placeholder in new_text:
                print(f"      Found placeholder: {placeholder}")
                
                # Cari nilai dari berbagai kemungkinan kolom
                replacement_value = ""
                
                for col_name in possible_columns:
                    if col_name in row and pd.notna(row[col_name]):
                        replacement_value = str(row[col_name]).strip()
                        print(f"        ✓ Using column '{col_name}': {replacement_value[:50]}")
                        break
                
                # Jika tidak ditemukan, cari dengan case-insensitive
                if not replacement_value:
                    # Cari dengan pattern matching di semua kolom
                    for col in row.index:
                        if str(col).lower() in [c.lower() for c in possible_columns]:
                            if pd.notna(row[col]):
                                replacement_value = str(row[col]).strip()
                                print(f"        ⚠ Found with case-insensitive '{col}': {replacement_value[:50]}")
                                break
                
                # Jika masih kosong, beri default value berdasarkan placeholder
                if not replacement_value:
                    if placeholder == '{{nik}}':
                        replacement_value = "N/A"
                        print(f"        ⚠ NIK not found, using default: N/A")
                    elif placeholder == '{{nama}}':
                        replacement_value = "N/A"
                        print(f"        ⚠ NAMA not found, using default: N/A")
                    else:
                        replacement_value = ""
                        print(f"        ⚠ Column not found for {placeholder}")
                
                # Replace text
                new_text = new_text.replace(placeholder, replacement_value)
                print(f"        Replaced: '{placeholder}' -> '{replacement_value[:30]}...'")
        
        # Set new text after all replacements
        if new_text != original_text:
            text_frame.text = new_text
            
            # Apply formatting
            try:
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Set font size berdasarkan konten
                        if "{{nama}}" in original_text or "Nama" in run.text:
                            run.font.size = Pt(15)
                            run.font.bold = True
                        elif "{{nik}}" in original_text or "NIK" in run.text:
                            run.font.size = Pt(15)
                            run.font.bold = True
                        elif "{{jabatan terakhir}}" in original_text or "Jabatan" in run.text:
                            run.font.size = Pt(15)
                            run.font.bold = False
                        else:
                            run.font.size = Pt(10.5)
                            run.font.bold = False
            except Exception as e:
                print(f"        ⚠ Formatting error: {e}")

def handle_jabatan_placeholder(text_frame, row):
    """
    Special handler untuk placeholder jabatan terakhir yang mungkin terpisah
    """
    original_text = text_frame.text
    
    # Normalize whitespace
    normalized_text = " ".join(original_text.split())
    
    # Patterns untuk mencari placeholder jabatan
    patterns = [
        r"{{jabatan terakhir}}",
        r"{{j\s*abatan\s*terakhir\s*}}",
        r"{{.*j.*abatan.*terakhir.*}}"
    ]
    
    for pattern in patterns:
        if re.search(pattern, normalized_text, re.IGNORECASE):
            jabatan = str(row.get("jabatan terakhir", ""))
            
            # Replace text
            text_frame.text = jabatan
            
            # Apply formatting
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(15)
                    run.font.bold = False
            
            return True
    
    return False

def validate_template(template_path: str) -> bool:
    """
    Validate template PowerPoint
    """
    try:
        prs = Presentation(template_path)
        
        if len(prs.slides) == 0:
            print("❌ Template tidak memiliki slide")
            return False
        
        # Check for placeholders
        slide = prs.slides[0]
        text_content = ""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content += shape.text_frame.text
        
        required_placeholders = [
            '{{nama}}', '{{nik}}', '{{executive summary}}',
            '{{education}}', '{{competency}}', '{{experience}}'
        ]
        
        missing = []
        for placeholder in required_placeholders:
            if placeholder not in text_content:
                missing.append(placeholder)
        
        if missing:
            print(f"⚠️ Warning: Missing placeholders: {', '.join(missing)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Template validation error: {e}")
        return False