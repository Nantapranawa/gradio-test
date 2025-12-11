import pandas as pd
from pptx import Presentation
from pptx.util import Pt
import os
import re

def create_formatted_resume_advanced():
    # Baca data
    df = pd.read_csv("D:/Project OCR Telkom/Result/hasil_analisis_terintegrasi_20251205_155730.csv", encoding='utf-8-sig')
    
    template_path = "D:/Project OCR Telkom/Template Talent Resume.pptx"
    output_dir = "output_advanced"
    os.makedirs(output_dir, exist_ok=True)
    
    for index, row in df.iterrows():
        print(f"Processing: {row['nama']}")
        
        # Buka template
        prs = Presentation(template_path)
        slide = prs.slides[0]
        
        # Debug: Cetak semua teks di slide
        print(f"\nDebug - All text frames for {row['nama']}:")
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"  Shape {i}: '{text[:100]}...'")
        
        # Cari dan ganti dengan formatting yang tepat
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            text_frame = shape.text_frame
            original_text = text_frame.text
            
            # Debug khusus untuk jabatan
            if any(keyword in original_text.lower() for keyword in ["jabatan", "terakhir", "{{j"]):
                print(f"\nFound potential jabatan text frame: '{original_text[:100]}...'")
            
            # 1. Nama dan NIK
            if "{{nama}}" in original_text and "{{nik}}" in original_text:
                new_text = original_text.replace("{{nama}}", str(row['nama']))
                new_text = new_text.replace("{{nik}}", str(row['nik']))
                text_frame.text = new_text
                
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(15)
                        run.font.bold = True
            
            # 2. Executive Summary
            elif "{{executive summary}}" in original_text:
                text_frame.text = original_text.replace(
                    "{{executive summary}}", 
                    str(row.get("summary executive", ""))
                )
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10.5)
                        run.font.bold = False
            
            # 3. Education
            elif "{{education}}" in original_text:
                text_frame.text = original_text.replace(
                    "{{education}}", 
                    str(row.get("education", ""))
                )
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10.5)
                        run.font.bold = False
            
            # 4. Jabatan Terakhir - PERBAIKAN UTAMA
            elif any(keyword in original_text for keyword in ["{{jabatan terakhir}}"]):
                # Cek jika ini adalah placeholder jabatan terakhir
                # Normalisasi whitespace (hapus line breaks, extra spaces)
                normalized_text = " ".join(original_text.split())
                
                print(f"  Normalized text: '{normalized_text}'")
                
                # Pattern untuk mencari placeholder dengan format terpisah
                pattern1 = r"{{jabatan terakhir}}"
                pattern2 = r"{{j\s*abatan\s*terakhir\s*}}"
                pattern3 = r"{{.*j.*abatan.*terakhir.*}}"
                
                # Debug patterns
                print(f"  Pattern1 match: {bool(re.search(pattern1, normalized_text, re.IGNORECASE))}")
                print(f"  Pattern2 match: {bool(re.search(pattern2, normalized_text, re.IGNORECASE))}")
                print(f"  Pattern3 match: {bool(re.search(pattern3, normalized_text, re.IGNORECASE))}")
                
                # Jika ditemukan placeholder untuk jabatan terakhir
                if (re.search(pattern1, normalized_text, re.IGNORECASE) or 
                    re.search(pattern2, normalized_text, re.IGNORECASE) or 
                    re.search(pattern3, normalized_text, re.IGNORECASE)):
                    
                    jabatan = str(row.get("jabatan terakhir", ""))
                    print(f"  Replacing with: '{jabatan}'")
                    
                    text_frame.text = jabatan
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(15)
                            run.font.bold = False
            
            # 5. Competency
            elif "{{competency}}" in original_text:
                text_frame.text = original_text.replace(
                    "{{competency}}", 
                    str(row.get("competency", ""))
                )
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10.5)
                        run.font.bold = False
            
            # 6. Experience
            elif "{{experience}}" in original_text:
                text_frame.text = original_text.replace(
                    "{{experience}}", 
                    str(row.get("experience", ""))
                )
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.bold = False
            
            # 7. Business Impact
            elif "{{business impact}}" in original_text:
                text_frame.text = original_text.replace(
                    "{{business impact}}", 
                    str(row.get("business impact", ""))
                )
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.bold = False
        
        # Simpan file
        output_path = os.path.join(output_dir, f"Resume_{row['nama'].replace(' ', '_')}.pptx")
        prs.save(output_path)
        print(f"  Saved: {output_path}")
    
    print(f"\nDone! Files saved in {output_dir}")

# Jalankan fungsi
create_formatted_resume_advanced()